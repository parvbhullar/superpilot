import io
import json
import os
import re
import zipfile
from collections.abc import Iterator
from email.parser import Parser as EmailParser
from pathlib import Path
from typing import Any
from typing import IO

import chardet
import docx  # type: ignore
import openpyxl  # type: ignore
import pptx  # type: ignore
from pypdf import PdfReader
from pypdf.errors import PdfStreamError

#from super_store.configs.constants import UNPOD_METADATA_FILENAME
#from super_store.file_processing.html_utils import parse_html_page_basic
#from super_store.utils.logger import setup_logger
from superpilot.core.logging.logging import setup_logger

logger = setup_logger()


TEXT_SECTION_SEPARATOR = "\n\n"


PLAIN_TEXT_FILE_EXTENSIONS = [
    ".txt",
    ".md",
    ".mdx",
    ".conf",
    ".log",
    ".json",
    ".xml",
    ".yml",
    ".yaml",
]

TABULAR_FILE_EXTENSIONS = [
    ".csv",
    ".tsv",
    ".record"
]


VALID_FILE_EXTENSIONS = PLAIN_TEXT_FILE_EXTENSIONS + [
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".eml",
    ".epub",
    ".html",
] + TABULAR_FILE_EXTENSIONS


def is_text_file_extension(file_name: str) -> bool:
    return any(file_name.endswith(ext) for ext in PLAIN_TEXT_FILE_EXTENSIONS)


def is_tabular_file_extension(file_name: str) -> bool:
    return any(file_name.endswith(ext) for ext in TABULAR_FILE_EXTENSIONS)


def get_file_ext(file_path_or_name: str | Path) -> str:
    _, extension = os.path.splitext(file_path_or_name)
    return extension


def check_file_ext_is_valid(ext: str) -> bool:
    return ext in VALID_FILE_EXTENSIONS


def detect_encoding(file: IO[bytes]) -> str:
    raw_data = file.read(50000)
    encoding = chardet.detect(raw_data)["encoding"] or "utf-8"
    file.seek(0)
    return encoding


def is_macos_resource_fork_file(file_name: str) -> bool:
    return os.path.basename(file_name).startswith("._") and file_name.startswith(
        "__MACOSX"
    )


# To include additional metadata in the search index, add a .unpod_metadata.json file
# to the zip file. This file should contain a list of objects with the following format:
# [{ "filename": "file1.txt", "link": "https://example.com/file1.txt" }]
def load_files_from_zip(
    zip_file_io: IO,
    ignore_macos_resource_fork_files: bool = True,
    ignore_dirs: bool = True,
) -> Iterator[tuple[zipfile.ZipInfo, IO[Any], dict[str, Any]]]:
    with zipfile.ZipFile(zip_file_io, "r") as zip_file:
        zip_metadata = {}
        try:
            metadata_file_info = zip_file.getinfo(UNPOD_METADATA_FILENAME)
            with zip_file.open(metadata_file_info, "r") as metadata_file:
                try:
                    zip_metadata = json.load(metadata_file)
                    if isinstance(zip_metadata, list):
                        # convert list of dicts to dict of dicts
                        zip_metadata = {d["filename"]: d for d in zip_metadata}
                except json.JSONDecodeError:
                    logger.warn(f"Unable to load {UNPOD_METADATA_FILENAME}")
        except KeyError:
            logger.info(f"No {UNPOD_METADATA_FILENAME} file")

        for file_info in zip_file.infolist():
            with zip_file.open(file_info.filename, "r") as file:
                if ignore_dirs and file_info.is_dir():
                    continue

                if (
                    ignore_macos_resource_fork_files
                    and is_macos_resource_fork_file(file_info.filename)
                ) or file_info.filename == UNPOD_METADATA_FILENAME:
                    continue
                yield file_info, file, zip_metadata.get(file_info.filename, {})


def _extract_unpod_metadata(line: str) -> dict | None:
    html_comment_pattern = r"<!--\s*UNPOD_METADATA=\{(.*?)\}\s*-->"
    hashtag_pattern = r"#UNPOD_METADATA=\{(.*?)\}"

    html_comment_match = re.search(html_comment_pattern, line)
    hashtag_match = re.search(hashtag_pattern, line)

    if html_comment_match:
        json_str = html_comment_match.group(1)
    elif hashtag_match:
        json_str = hashtag_match.group(1)
    else:
        return None

    try:
        return json.loads("{" + json_str + "}")
    except json.JSONDecodeError:
        return None


def read_text_file(
    file: IO,
    encoding: str = "utf-8",
    errors: str = "replace",
    ignore_unpod_metadata: bool = True,
) -> tuple[str, dict]:
    metadata = {}
    file_content_raw = ""
    for ind, line in enumerate(file):
        try:
            line = line.decode(encoding) if isinstance(line, bytes) else line
        except UnicodeDecodeError:
            line = (
                line.decode(encoding, errors=errors)
                if isinstance(line, bytes)
                else line
            )

        if ind == 0:
            metadata_or_none = (
                None if ignore_unpod_metadata else _extract_unpod_metadata(line)
            )
            if metadata_or_none is not None:
                metadata = metadata_or_none
            else:
                file_content_raw += line
        else:
            file_content_raw += line

    return file_content_raw, metadata


def pdf_to_text(file: IO[Any], pdf_pass: str | None = None) -> str:
    try:
        pdf_reader = PdfReader(file)

        # If marked as encrypted and a password is provided, try to decrypt
        if pdf_reader.is_encrypted and pdf_pass is not None:
            decrypt_success = False
            if pdf_pass is not None:
                try:
                    decrypt_success = pdf_reader.decrypt(pdf_pass) != 0
                except Exception:
                    logger.error("Unable to decrypt pdf")
            else:
                logger.info("No Password available to to decrypt pdf")

            if not decrypt_success:
                # By user request, keep files that are unreadable just so they
                # can be discoverable by title.
                return ""

        return TEXT_SECTION_SEPARATOR.join(
            page.extract_text() for page in pdf_reader.pages
        )
    except PdfStreamError:
        logger.exception("PDF file is not a valid PDF")
    except Exception:
        logger.exception("Failed to read PDF")

    # File is still discoverable by title
    # but the contents are not included as they cannot be parsed
    return ""


def docx_to_text(file: IO[Any]) -> str:
    doc = docx.Document(file)
    full_text = [para.text for para in doc.paragraphs]
    return TEXT_SECTION_SEPARATOR.join(full_text)


def pptx_to_text(file: IO[Any]) -> str:
    presentation = pptx.Presentation(file)
    text_content = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        extracted_text = f"\nSlide {slide_number}:\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                extracted_text += shape.text + "\n"
        text_content.append(extracted_text)
    return TEXT_SECTION_SEPARATOR.join(text_content)


def xlsx_to_text(file: IO[Any]) -> str:
    workbook = openpyxl.load_workbook(file)
    text_content = []
    for sheet in workbook.worksheets:
        sheet_string = "\n".join(
            ",".join(map(str, row))
            for row in sheet.iter_rows(min_row=1, values_only=True)
        )
        text_content.append(sheet_string)
    return TEXT_SECTION_SEPARATOR.join(text_content)


def eml_to_text(file: IO[Any]) -> str:
    text_file = io.TextIOWrapper(file, encoding=detect_encoding(file))
    parser = EmailParser()
    message = parser.parse(text_file)
    text_content = []
    for part in message.walk():
        if part.get_content_type().startswith("text/plain"):
            text_content.append(part.get_payload())
    return TEXT_SECTION_SEPARATOR.join(text_content)


def epub_to_text(file: IO[Any]) -> str:
    with zipfile.ZipFile(file) as epub:
        text_content = []
        for item in epub.infolist():
            if item.filename.endswith(".xhtml") or item.filename.endswith(".html"):
                with epub.open(item) as html_file:
                    text_content.append(parse_html_page_basic(html_file))
        return TEXT_SECTION_SEPARATOR.join(text_content)


def file_io_to_text(file: IO[Any]) -> str:
    encoding = detect_encoding(file)
    file_content_raw, _ = read_text_file(file, encoding=encoding)
    return file_content_raw


def extract_file_text(
    file_name: str | None,
    file: IO[Any],
    break_on_unprocessable: bool = True,
) -> str:
    if not file_name:
        return file_io_to_text(file)

    extension = get_file_ext(file_name)
    if not check_file_ext_is_valid(extension):
        if break_on_unprocessable:
            raise RuntimeError(f"Unprocessable file type: {file_name}")
        else:
            logger.warning(f"Unprocessable file type: {file_name}")
            return ""

    if extension == ".pdf":
        return pdf_to_text(file=file)

    elif extension == ".docx":
        return docx_to_text(file)

    elif extension == ".pptx":
        return pptx_to_text(file)

    elif extension == ".xlsx":
        return xlsx_to_text(file)

    elif extension == ".eml":
        return eml_to_text(file)

    elif extension == ".epub":
        return epub_to_text(file)

    elif extension == ".html":
        return parse_html_page_basic(file)

    else:
        return file_io_to_text(file)
